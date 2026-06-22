import io
from pptx import Presentation
from pptx.util import Inches
from processors.pptx import extract_texts, reinsert_texts


def _make_pptx(texts: list[str]) -> bytes:
    prs = Presentation()
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    for text in texts:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx_with_table(rows: list[list[str]]) -> bytes:
    prs = Presentation()
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    row_count = len(rows)
    col_count = len(rows[0]) if rows else 1
    table = slide.shapes.add_table(row_count, col_count, Inches(1), Inches(1), Inches(6), Inches(3)).table
    for ri, row in enumerate(rows):
        for ci, text in enumerate(row):
            table.cell(ri, ci).text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_texts():
    data = _make_pptx(["안녕하세요", "테스트"])
    segments = extract_texts(data)
    texts = [s["text"] for s in segments]
    assert "안녕하세요" in texts
    assert "테스트" in texts


def test_reinsert_texts():
    data = _make_pptx(["안녕하세요"])
    segments = extract_texts(data)
    result = reinsert_texts(data, segments, ["Hello"])
    prs = Presentation(io.BytesIO(result))
    all_texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_texts.append(shape.text_frame.text)
    assert "Hello" in all_texts


def test_extract_texts_table():
    data = _make_pptx_with_table([["부서", "주요 요구사항"], ["SCM팀", "데이터 매칭"]])
    segments = extract_texts(data)
    texts = [s["text"] for s in segments]
    assert "부서" in texts
    assert "주요 요구사항" in texts
    assert "SCM팀" in texts
    assert "데이터 매칭" in texts


def test_reinsert_texts_table():
    data = _make_pptx_with_table([["부서", "주요 요구사항"], ["SCM팀", "데이터 매칭"]])
    segments = extract_texts(data)
    translated = ["Department", "Key Requirements", "SCM Team", "Data Matching"]
    result = reinsert_texts(data, segments, translated)
    prs = Presentation(io.BytesIO(result))
    all_texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        all_texts.append(cell.text_frame.text)
    assert "Department" in all_texts
    assert "Key Requirements" in all_texts
    assert "SCM Team" in all_texts
    assert "Data Matching" in all_texts
