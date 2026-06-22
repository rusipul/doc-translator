import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def extract_texts(file_bytes: bytes) -> list[dict]:
    prs = Presentation(io.BytesIO(file_bytes))
    segments = []
    for si, slide in enumerate(prs.slides):
        for shape in _iter_shapes(slide.shapes):
            shape_id = shape.shape_id
            if shape.has_text_frame:
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    for ri, run in enumerate(para.runs):
                        if run.text.strip():
                            segments.append({
                                "text": run.text,
                                "key": (si, shape_id, -1, -1, pi, ri),
                            })
            elif shape.has_table:
                for rowi, row in enumerate(shape.table.rows):
                    for coli, cell in enumerate(row.cells):
                        for pi, para in enumerate(cell.text_frame.paragraphs):
                            for ri, run in enumerate(para.runs):
                                if run.text.strip():
                                    segments.append({
                                        "text": run.text,
                                        "key": (si, shape_id, rowi, coli, pi, ri),
                                    })
    return segments


def reinsert_texts(file_bytes: bytes, segments: list[dict], translated: list[str]) -> bytes:
    prs = Presentation(io.BytesIO(file_bytes))
    key_to_translation = {
        tuple(s["key"]): t for s, t in zip(segments, translated)
    }
    for si, slide in enumerate(prs.slides):
        for shape in _iter_shapes(slide.shapes):
            shape_id = shape.shape_id
            if shape.has_text_frame:
                for pi, para in enumerate(shape.text_frame.paragraphs):
                    for ri, run in enumerate(para.runs):
                        key = (si, shape_id, -1, -1, pi, ri)
                        if key in key_to_translation:
                            run.text = key_to_translation[key]
            elif shape.has_table:
                for rowi, row in enumerate(shape.table.rows):
                    for coli, cell in enumerate(row.cells):
                        for pi, para in enumerate(cell.text_frame.paragraphs):
                            for ri, run in enumerate(para.runs):
                                key = (si, shape_id, rowi, coli, pi, ri)
                                if key in key_to_translation:
                                    run.text = key_to_translation[key]
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
